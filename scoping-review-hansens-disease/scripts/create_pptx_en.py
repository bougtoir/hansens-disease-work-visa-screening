#!/usr/bin/env python3
"""Generate PPTX with figures/tables (English) - 1 figure/table per slide.
Code-generated charts -> embedded as images.
PRISMA flow diagram -> editable PowerPoint shapes.
Tables -> editable PowerPoint tables.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

FIGURES_DIR = "/home/ubuntu/scoping_review/figures"
OUTPUT_DIR = "/home/ubuntu/scoping_review/docx"
os.makedirs(OUTPUT_DIR, exist_ok=True)

DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE_BG = RGBColor(0xBB, 0xDE, 0xFB)
GREEN_BG = RGBColor(0xC8, 0xE6, 0xC9)
LIGHT_RED_BG = RGBColor(0xFF, 0xCD, 0xD2)
YELLOW_BG = RGBColor(0xFF, 0xF9, 0xC4)
DARK_GREEN = RGBColor(0x2E, 0x7D, 0x32)
DARK_RED = RGBColor(0xC6, 0x28, 0x28)
BLUE_BORDER = RGBColor(0x15, 0x65, 0xC0)
AMBER_BORDER = RGBColor(0xF9, 0xA8, 0x25)
PINK_BG = RGBColor(0xFC, 0xE4, 0xEC)
LIGHT_GREEN_BG2 = RGBColor(0xE8, 0xF5, 0xE9)


def set_cell_format(cell, text, font_size=10, bold=False, color=None, fill=None, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Hansen's Disease Screening in Work Visa Medical Examinations"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "A Global Scoping Review of Administrative Requirements\nThat Diverge from Standard Medical Practice"
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7.0), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Figures and Tables\nPRISMA-ScR Scoping Review"
    run3.font.size = Pt(16)
    run3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    txBox3 = slide.shapes.add_textbox(Inches(2.0), Inches(6.5), Inches(6.0), Inches(0.5))
    tf3 = txBox3.text_frame
    p4 = tf3.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = "[Author Names] | 2026"
    run4.font.size = Pt(12)
    run4.font.italic = True
    run4.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def add_image_slide(prs, image_path, title_text, caption_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    if os.path.exists(image_path):
        pic = slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.0), Inches(9.0))
        pw, ph = pic.width, pic.height
        ratio = min(Inches(9.0) / pw, Inches(5.5) / ph)
        pic.width = int(pw * ratio)
        pic.height = int(ph * ratio)
        pic.left = int((Inches(10.0) - pic.width) / 2)
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.7))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = caption_text
    run2.font.size = Pt(9)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def _add_arrow_xml(slide, x1, y1, x2, y2, color_hex="1565C0"):
    spTree = slide.shapes._spTree
    shape_id = len(spTree) + 100
    cx = x2 - x1
    cy = y2 - y1
    flip_h = ""
    flip_v = ""
    off_x, off_y = x1, y1
    if cx < 0:
        flip_h = ' flipH="1"'
        off_x = x2
        cx = abs(cx)
    if cy < 0:
        flip_v = ' flipV="1"'
        off_y = y2
        cy = abs(cy)
    xml_str = (
        '<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<p:nvCxnSpPr>'
        f'<p:cNvPr id="{shape_id}" name="Arrow {shape_id}"/>'
        '<p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>'
        f'<p:spPr><a:xfrm{flip_h}{flip_v}>'
        f'<a:off x="{int(off_x)}" y="{int(off_y)}"/>'
        f'<a:ext cx="{int(cx)}" cy="{int(cy)}"/>'
        '</a:xfrm>'
        '<a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>'
        f'<a:ln w="19050"><a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
        '<a:tailEnd type="triangle" w="med" len="med"/></a:ln>'
        '</p:spPr></p:cxnSp>'
    )
    spTree.append(etree.fromstring(xml_str))


def add_editable_prisma_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(9.4), Inches(0.45))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Figure 4: PRISMA-ScR Flow Diagram"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    def add_box(left, top, width, height, text, fill_color, border_color, font_size=9):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
        return shape

    def add_phase_label(left, top, width, height, text):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
        shape.line.color.rgb = BLUE_BORDER
        shape.line.width = Pt(1.0)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = BLUE_BORDER

    def arrow_b2t(s_from, s_to):
        _add_arrow_xml(slide,
            s_from.left + s_from.width // 2, s_from.top + s_from.height,
            s_to.left + s_to.width // 2, s_to.top)

    def arrow_r2l(s_from, s_to):
        _add_arrow_xml(slide,
            s_from.left + s_from.width, s_from.top + s_from.height // 2,
            s_to.left, s_to.top + s_to.height // 2)

    add_phase_label(Inches(0.1), Inches(0.6), Inches(0.85), Inches(0.5), "IDENTIFICATION")
    add_phase_label(Inches(0.1), Inches(1.55), Inches(0.85), Inches(0.5), "SCREENING")
    add_phase_label(Inches(0.1), Inches(2.7), Inches(0.85), Inches(0.5), "ELIGIBILITY")
    add_phase_label(Inches(0.1), Inches(4.2), Inches(0.85), Inches(0.5), "INCLUDED")

    box1 = add_box(Inches(1.3), Inches(0.55), Inches(6.5), Inches(0.6),
        "Countries/territories identified for review\n(n = 197: 193 UN member states + 4 observer entities)",
        LIGHT_BLUE_BG, BLUE_BORDER, 10)
    box2 = add_box(Inches(1.3), Inches(1.5), Inches(6.5), Inches(0.65),
        "Sources screened per country:\nGovernment websites, legislation, official forms,\nILEP database, IOM reports, academic literature",
        LIGHT_BLUE_BG, BLUE_BORDER, 9)
    box3 = add_box(Inches(1.3), Inches(2.55), Inches(5.2), Inches(0.6),
        "English-language sources: characterized\n(n = 110: 58 disease-specific + 52 no requirement)",
        GREEN_BG, DARK_GREEN, 9)
    box3_excl = add_box(Inches(7.1), Inches(2.55), Inches(2.5), Inches(0.6),
        "English insufficient\n(n = 87)", LIGHT_RED_BG, DARK_RED, 9)
    box4 = add_box(Inches(1.3), Inches(3.55), Inches(5.2), Inches(0.6),
        "Supplementary multilingual research:\n5 countries confirmed + 82 unresolved",
        GREEN_BG, DARK_GREEN, 9)
    box4_excl = add_box(Inches(7.1), Inches(3.55), Inches(2.5), Inches(0.6),
        "Ultimately unresolved\ndespite multilingual search\n(n = 82)", YELLOW_BG, AMBER_BORDER, 8)
    box5a = add_box(Inches(1.1), Inches(4.65), Inches(3.6), Inches(0.75),
        "Countries with Hansen's disease\nexplicitly named in work visa\nmedical requirements (n = 20)",
        LIGHT_RED_BG, DARK_RED, 9)
    box5b = add_box(Inches(5.3), Inches(4.65), Inches(3.6), Inches(0.75),
        "Countries with disease-specific\nscreening but NO Hansen's\ndisease (n = 38)",
        LIGHT_GREEN_BG2, DARK_GREEN, 9)
    box6a = add_box(Inches(1.1), Inches(5.8), Inches(3.6), Inches(1.3),
        "By region:\nGCC/Middle East: 6\nSE/East Asia: 5\nAmericas: 3\nAfrica: 2\nEurope: 2\nSouth Asia: 2",
        PINK_BG, DARK_RED, 8)
    box6b = add_box(Inches(5.3), Inches(5.8), Inches(3.6), Inches(1.3),
        "Top screened diseases (excl. Hansen's):\nTB: 38 countries\nHIV: 35 countries\nSyphilis: 28 countries\nHepatitis B: 22 countries",
        LIGHT_GREEN_BG2, DARK_GREEN, 8)

    arrow_b2t(box1, box2)
    arrow_b2t(box2, box3)
    arrow_r2l(box3, box3_excl)
    arrow_b2t(box3, box4)
    arrow_r2l(box4, box4_excl)
    x_mid = box4.left + box4.width // 2
    y_bot4 = box4.top + box4.height
    _add_arrow_xml(slide, x_mid, y_bot4, box5a.left + box5a.width // 2, box5a.top)
    _add_arrow_xml(slide, x_mid, y_bot4, box5b.left + box5b.width // 2, box5b.top)
    arrow_b2t(box5a, box6a)
    arrow_b2t(box5b, box6b)


def add_table_slide(prs, title_text, headers, data, col_widths, header_fill):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    rows = len(data) + 1
    cols = len(headers)
    tbl = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(0.7), Inches(9.4), Inches(6.5)).table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)
    for i, h in enumerate(headers):
        set_cell_format(tbl.cell(0, i), h, 8, True, WHITE, header_fill, PP_ALIGN.CENTER)
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            set_cell_format(tbl.cell(r+1, c), val, 7, False, alignment=PP_ALIGN.LEFT)
            if r % 2 == 1:
                tbl.cell(r+1, c).fill.solid()
                tbl.cell(r+1, c).fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)


def create_pptx_en():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig1_world_map.png"),
        "Figure 1: Global Distribution of Hansen's Disease Screening",
        "Schematic regional map showing countries that screen for Hansen's disease (red), "
        "disease-specific screening without Hansen's (blue), and no disease-specific requirements (green).")

    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig2_disease_bar.png"),
        "Figure 2: Diseases Named in Work Visa Medical Screening Requirements",
        "Horizontal bar chart across 58 countries. Hansen's disease (red) is the 5th most commonly screened (34.5%).")

    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig3_regional_donut.png"),
        "Figure 3: Regional Distribution and Legal Provisions",
        "(A) Regional distribution of 20 countries screening for Hansen's disease. "
        "(B) Nature of Hansen's disease provisions in immigration law.")

    add_editable_prisma_flow(prs)

    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig5_transmissibility.png"),
        "Figure 5: Disease Transmissibility vs. Screening Frequency",
        "Hansen's disease (red) shows disproportionately high screening despite very low transmissibility. "
        "Bubble size proportional to number of countries screening for each disease.")

    # Slide: Figure 6 - Sankey Diagram (NEW)
    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig6_sankey_en.png"),
        "Figure 6: Data Accessibility Flow Diagram",
        "Two-stage search strategy: 197 countries → 110 via English sources (58 disease-specific + "
        "52 no requirement), 5 via multilingual research, 82 unreachable. "
        "Among 63 with confirmed screening, 20 named Hansen's disease.")

    # Slide: Figure 7 - Accessibility Map (NEW)
    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig7_accessibility_en.png"),
        "Figure 7: Geographic Distribution of Data Accessibility",
        "Green: reached via English sources (110, 55.8%); Orange: reached via multilingual research "
        "(5, 2.5%); Red: unreachable (82, 41.6%). Gray: not in study scope.")

    t1h = ['Region', 'Country', 'Legal Instrument', 'Provision Type', 'Year']
    t1d = [
        ['GCC/Middle East', 'Saudi Arabia', 'GAMCA/WAFID Regulations 6th Ed.', 'Automatic exclusion', '2021'],
        ['', 'UAE', 'GAMCA/WAFID Regulations', 'Automatic exclusion', '2021'],
        ['', 'Qatar', 'GAMCA/WAFID Regulations', 'Automatic exclusion', '2021'],
        ['', 'Kuwait', 'GAMCA/WAFID Regulations', 'Automatic exclusion', '2021'],
        ['', 'Oman', 'GAMCA/WAFID Regulations', 'Automatic exclusion', '2021'],
        ['', 'Bahrain', 'GAMCA/WAFID Regulations', 'Automatic exclusion', '2021'],
        ['East/SE Asia', 'China', 'Foreigner Physical Exam. Form', 'Certification of absence', '\u2014'],
        ['', 'Thailand', 'Emergency Decree B.E.2560', 'Prohibited disease', '2017'],
        ['', 'Taiwan', 'Work Permit Regulations', 'Certification of absence', '\u2014'],
        ['', 'Malaysia', 'FOMEMA Requirements', 'Medical unfitness', '\u2014'],
        ['', 'Philippines', 'OWWA Requirements', 'Certification of absence', '\u2014'],
        ['Africa', 'South Africa', 'Immigration Act 2002; BI-811', 'Certification of absence', '2002'],
        ['', 'Namibia', 'Immigration Control Act', 'Medical prohibition', '1993'],
        ['Europe', 'Russia', 'Federal Law No.115-FZ', 'Deportation ground', '2002'],
        ['', 'Malta', 'Immigration Act Cap.217', 'Medical inadmissibility', '\u2014'],
        ['Americas', 'United States', 'INA \u00a7212(a)(1)(A)(i)', 'Class A condition*', '1952\u2020'],
        ['', 'Barbados', 'Immigration Act Cap.190', 'Medical prohibition', '\u2014'],
        ['', 'US Virgin Islands', 'US Federal Law', 'Class A condition*', '\u2014'],
        ['South Asia', 'India', 'State-level Employment Acts', 'Employment restriction', 'Various'],
        ['', 'Hong Kong SAR', 'Immigration Ordinance Cap.115', 'Medical inadmissibility', '\u2014'],
    ]
    add_table_slide(prs,
        "Table 1: Countries with Hansen's Disease in Work Visa Medical Requirements (n=20)",
        t1h, t1d, [1.3, 1.5, 2.8, 2.0, 0.8], DARK_BLUE)

    t2h = ['Country', 'Diseases Screened', 'Notable Features', "Hansen's Disease"]
    t2d = [
        ['Canada', 'TB, syphilis', 'Immigration Medical Exam (IME)', 'Not included'],
        ['Australia', 'TB, HIV, Hep B/C', 'Panel physician system (Bupa)', 'Not included'],
        ['United Kingdom', 'TB (visas >6 months)', 'IOM panel physician screening', 'Not included'],
        ['New Zealand', 'TB, HIV, Hep B', 'Immigration NZ panel physicians', 'Not included'],
        ['Japan', 'TB only (select nationalities)', 'JPETS pre-entry TB screening', 'Not included'],
        ['South Korea', 'TB, HIV, drugs', 'E-9 visa medical examination', 'Not included'],
        ['Singapore', 'TB, HIV, syphilis, Hep B, malaria', 'Foreign Worker Medical Exam', 'Not included'],
        ['Israel', 'TB, HIV, Hep B/C', 'Entry medical requirements', 'Not included'],
        ['Kenya', 'TB, HIV', 'Work permit medical certificate', 'Not included'],
        ['Nigeria', 'TB, HIV, Hep B', 'Expatriate quota medical exam', 'Not included'],
        ['Ghana', 'TB, HIV, yellow fever', 'Immigration medical certificate', 'Not included'],
        ['Brazil', 'Yellow fever (vaccination)', 'No work visa medical exam', 'Not included'],
        ['Mexico', 'None specified', 'No standard medical exam', 'Not included'],
        ['Argentina', 'None specified', 'No disease-specific exam', 'Not included'],
        ['Chile', 'TB, HIV (select visas)', 'Temporary resident visa medical', 'Not included'],
    ]
    add_table_slide(prs,
        "Table 2: Countries with Disease-Specific Screening but NOT Including Hansen's Disease (n=15)",
        t2h, t2d, [1.5, 2.5, 3.0, 1.4], DARK_GREEN)

    out = os.path.join(OUTPUT_DIR, "Scoping_Review_Figures_EN.pptx")
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == '__main__':
    create_pptx_en()
