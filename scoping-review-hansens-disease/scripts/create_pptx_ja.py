#!/usr/bin/env python3
"""Generate PPTX with figures/tables (Japanese) - 1 figure/table per slide.
Code-generated charts -> embedded as images.
PRISMA flow diagram -> editable PowerPoint shapes.
Tables -> editable PowerPoint tables.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

FIGURES_DIR = "/home/ubuntu/scoping_review/figures"
OUTPUT_DIR = "/home/ubuntu/scoping_review/docx"
os.makedirs(OUTPUT_DIR, exist_ok=True)

DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE_BG = RGBColor(0xBB, 0xDE, 0xFB)
GREEN_BG = RGBColor(0xC8, 0xE6, 0xC9)
LIGHT_RED_BG = RGBColor(0xFF, 0xCD, 0xD2)
YELLOW_BG = RGBColor(0xFF, 0xF9, 0xC4)
DARK_GREEN = RGBColor(0x2E, 0x7D, 0x32)
DARK_RED = RGBColor(0xC6, 0x28, 0x28)
BLUE_BORDER = RGBColor(0x15, 0x65, 0xC0)
AMBER_BORDER = RGBColor(0xF9, 0xA8, 0x25)
PINK_BG = RGBColor(0xFC, 0xE4, 0xEC)
LIGHT_GREEN_BG2 = RGBColor(0xE8, 0xF5, 0xE9)


def set_cell_format(cell, text, font_size=10, bold=False, color=None, fill=None, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "労働ビザ医療検査における\nハンセン病スクリーニング"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "標準医療と乖離した行政要求に関する\nグローバルスコーピングレビュー"
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7.0), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "図表一覧\nPRISMA-ScR スコーピングレビュー"
    run3.font.size = Pt(16)
    run3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    txBox3 = slide.shapes.add_textbox(Inches(2.0), Inches(6.5), Inches(6.0), Inches(0.5))
    tf3 = txBox3.text_frame
    p4 = tf3.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = "[著者名] | 2026"
    run4.font.size = Pt(12)
    run4.font.italic = True
    run4.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def add_image_slide(prs, image_path, title_text, caption_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    if os.path.exists(image_path):
        pic = slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.0), Inches(9.0))
        pw, ph = pic.width, pic.height
        ratio = min(Inches(9.0) / pw, Inches(5.5) / ph)
        pic.width = int(pw * ratio)
        pic.height = int(ph * ratio)
        pic.left = int((Inches(10.0) - pic.width) / 2)
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.7))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = caption_text
    run2.font.size = Pt(9)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def _add_arrow_xml(slide, x1, y1, x2, y2, color_hex="1565C0"):
    spTree = slide.shapes._spTree
    shape_id = len(spTree) + 100
    cx = x2 - x1
    cy = y2 - y1
    flip_h = ""
    flip_v = ""
    off_x, off_y = x1, y1
    if cx < 0:
        flip_h = ' flipH="1"'
        off_x = x2
        cx = abs(cx)
    if cy < 0:
        flip_v = ' flipV="1"'
        off_y = y2
        cy = abs(cy)
    xml_str = (
        '<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<p:nvCxnSpPr>'
        f'<p:cNvPr id="{shape_id}" name="Arrow {shape_id}"/>'
        '<p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>'
        f'<p:spPr><a:xfrm{flip_h}{flip_v}>'
        f'<a:off x="{int(off_x)}" y="{int(off_y)}"/>'
        f'<a:ext cx="{int(cx)}" cy="{int(cy)}"/>'
        '</a:xfrm>'
        '<a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>'
        f'<a:ln w="19050"><a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
        '<a:tailEnd type="triangle" w="med" len="med"/></a:ln>'
        '</p:spPr></p:cxnSp>'
    )
    spTree.append(etree.fromstring(xml_str))


def add_editable_prisma_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(9.4), Inches(0.45))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "図4：PRISMA-ScR フローダイアグラム"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    def add_box(left, top, width, height, text, fill_color, border_color, font_size=9):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
        return shape

    def add_phase_label(left, top, width, height, text):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
        shape.line.color.rgb = BLUE_BORDER
        shape.line.width = Pt(1.0)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = BLUE_BORDER

    def arrow_b2t(s_from, s_to):
        _add_arrow_xml(slide,
            s_from.left + s_from.width // 2, s_from.top + s_from.height,
            s_to.left + s_to.width // 2, s_to.top)

    def arrow_r2l(s_from, s_to):
        _add_arrow_xml(slide,
            s_from.left + s_from.width, s_from.top + s_from.height // 2,
            s_to.left, s_to.top + s_to.height // 2)

    # Phase labels (Japanese)
    add_phase_label(Inches(0.1), Inches(0.6), Inches(0.85), Inches(0.5), "同定\n(IDENTIFICATION)")
    add_phase_label(Inches(0.1), Inches(1.55), Inches(0.85), Inches(0.5), "スクリーニング\n(SCREENING)")
    add_phase_label(Inches(0.1), Inches(2.7), Inches(0.85), Inches(0.5), "適格性\n(ELIGIBILITY)")
    add_phase_label(Inches(0.1), Inches(4.2), Inches(0.85), Inches(0.5), "採択\n(INCLUDED)")

    box1 = add_box(Inches(1.3), Inches(0.55), Inches(6.5), Inches(0.6),
        "レビュー対象として同定された国・地域\n(n = 197：国連加盟国193 + オブザーバー等4)",
        LIGHT_BLUE_BG, BLUE_BORDER, 10)
    box2 = add_box(Inches(1.3), Inches(1.5), Inches(6.5), Inches(0.65),
        "各国につきスクリーニングした情報源：\n政府ウェブサイト、法令、公式書式、\nILEPデータベース、IOM報告書、学術文献",
        LIGHT_BLUE_BG, BLUE_BORDER, 9)
    box3 = add_box(Inches(1.3), Inches(2.55), Inches(5.2), Inches(0.6),
        "英語情報源で特定\n（n = 110：疾病特異的58 + 要件なし52）",
        GREEN_BG, DARK_GREEN, 9)
    box3_excl = add_box(Inches(7.1), Inches(2.55), Inches(2.5), Inches(0.6),
        "英語情報不十分\n（n = 87）", LIGHT_RED_BG, DARK_RED, 9)
    box4 = add_box(Inches(1.3), Inches(3.55), Inches(5.2), Inches(0.6),
        "補足的多言語調査：\n5カ国確認 + 82カ国未解決",
        GREEN_BG, DARK_GREEN, 9)
    box4_excl = add_box(Inches(7.1), Inches(3.55), Inches(2.5), Inches(0.6),
        "多言語調査後も\n最終的に未解決\n（n = 82）", YELLOW_BG, AMBER_BORDER, 8)
    box5a = add_box(Inches(1.1), Inches(4.65), Inches(3.6), Inches(0.75),
        "ハンセン病が労働ビザ医療要件に\n明示的に記載されている国\n（n = 20）",
        LIGHT_RED_BG, DARK_RED, 9)
    box5b = add_box(Inches(5.3), Inches(4.65), Inches(3.6), Inches(0.75),
        "疾病特異的スクリーニングあるが\nハンセン病を含まない国\n（n = 38）",
        LIGHT_GREEN_BG2, DARK_GREEN, 9)
    box6a = add_box(Inches(1.1), Inches(5.8), Inches(3.6), Inches(1.3),
        "地域別内訳：\nGCC/中東：6カ国\n東南・東アジア：5カ国・地域\n南北アメリカ：3カ国\nアフリカ：2カ国\n欧州：2カ国 | 南アジア：2カ国・地域",
        PINK_BG, DARK_RED, 8)
    box6b = add_box(Inches(5.3), Inches(5.8), Inches(3.6), Inches(1.3),
        "主要スクリーニング疾病（ハンセン病除く）：\n結核：38カ国\nHIV：35カ国\n梅毒：28カ国\nB型肝炎：22カ国",
        LIGHT_GREEN_BG2, DARK_GREEN, 8)

    arrow_b2t(box1, box2)
    arrow_b2t(box2, box3)
    arrow_r2l(box3, box3_excl)
    arrow_b2t(box3, box4)
    arrow_r2l(box4, box4_excl)
    x_mid = box4.left + box4.width // 2
    y_bot4 = box4.top + box4.height
    _add_arrow_xml(slide, x_mid, y_bot4, box5a.left + box5a.width // 2, box5a.top)
    _add_arrow_xml(slide, x_mid, y_bot4, box5b.left + box5b.width // 2, box5b.top)
    arrow_b2t(box5a, box6a)
    arrow_b2t(box5b, box6b)


def add_table_slide(prs, title_text, headers, data, col_widths, header_fill):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    rows = len(data) + 1
    cols = len(headers)
    tbl = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(0.7), Inches(9.4), Inches(6.5)).table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)
    for i, h in enumerate(headers):
        set_cell_format(tbl.cell(0, i), h, 8, True, WHITE, header_fill, PP_ALIGN.CENTER)
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            set_cell_format(tbl.cell(r+1, c), val, 7, False, alignment=PP_ALIGN.LEFT)
            if r % 2 == 1:
                tbl.cell(r+1, c).fill.solid()
                tbl.cell(r+1, c).fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)


def create_pptx_ja():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    # Slide 2: Figure 1 - World Map (image, English chart + Japanese caption)
    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig1_world_map.png"),
        "図1：労働ビザ医療検査におけるハンセン病スクリーニングの世界分布",
        "ハンセン病スクリーニング実施国（赤）、ハンセン病以外の疾病特異的スクリーニング（青）、"
        "疾病特異的要件なし（緑）を示す概略地域図。バブルサイズは国数に比例。")

    # Slide 3: Figure 2 - Disease Bar Chart (image)
    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig2_disease_bar.png"),
        "図2：労働ビザ医療スクリーニング要件で指定される疾病",
        "58カ国の疾病特異的スクリーニングにおける疾病頻度を示す横棒グラフ。"
        "ハンセン病（赤）は5番目に多くスクリーニングされる疾病（34.5%）。")

    # Slide 4: Figure 3 - Regional Donut (image)
    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig3_regional_donut.png"),
        "図3：地域分布と法的規定の類型",
        "(A) ハンセン病スクリーニングを行う20カ国・地域の地域分布。"
        "(B) 移民法におけるハンセン病規定の性質（自動排除、証明書、等）。")

    # Slide 5: Figure 4 - PRISMA Flow (EDITABLE shapes, Japanese text)
    add_editable_prisma_flow(prs)

    # Slide 6: Figure 5 - Transmissibility Scatter (image)
    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig5_transmissibility.png"),
        "図5：疾病の感染性とスクリーニング頻度の関係",
        "ハンセン病（赤）は極めて低い感染性（R\u2080）にもかかわらず不釣り合いに高いスクリーニング頻度を示す。"
        "バブルサイズは各疾病をスクリーニングする国数に比例。")

    # Slide: 図6 - Sankey Diagram (NEW)
    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig6_sankey_ja.png"),
        "図6：データ到達フロー図",
        "2段階検索戦略：197カ国 → 英語情報源で110カ国（疾病特異的58 + 要件なし52）、"
        "多言語調査で5カ国、到達不可82カ国。スクリーニング確認63カ国のうち20カ国がハンセン病を明記。")

    # Slide: 図7 - Accessibility Map (NEW)
    add_image_slide(prs, os.path.join(FIGURES_DIR, "fig7_accessibility_ja.png"),
        "図7：データ到達可能性の地理的分布",
        "緑：英語情報源で到達（110カ国、55.8%）、橙：多言語調査で到達（5カ国、2.5%）、"
        "赤：到達不可（82カ国、41.6%）。灰色：調査対象外。")

    # Slide: Table 1 - 20 countries (Japanese headers)
    t1h = ['地域', '国名', '法的根拠', '規定の種類', '年']
    t1d = [
        ['GCC/中東', 'サウジアラビア', 'GAMCA/WAFID規則 第6版', '自動排除', '2021'],
        ['', 'UAE', 'GAMCA/WAFID規則', '自動排除', '2021'],
        ['', 'カタール', 'GAMCA/WAFID規則', '自動排除', '2021'],
        ['', 'クウェート', 'GAMCA/WAFID規則', '自動排除', '2021'],
        ['', 'オマーン', 'GAMCA/WAFID規則', '自動排除', '2021'],
        ['', 'バーレーン', 'GAMCA/WAFID規則', '自動排除', '2021'],
        ['東・東南アジア', '中国', '外国人体格検査表', '非罹患証明', '\u2014'],
        ['', 'タイ', '緊急勅令B.E.2560', '禁止疾病', '2017'],
        ['', '台湾', '労働許可規則', '非罹患証明', '\u2014'],
        ['', 'マレーシア', 'FOMEMA要件', '医学的不適格', '\u2014'],
        ['', 'フィリピン', 'OWWA要件', '非罹患証明', '\u2014'],
        ['アフリカ', '南アフリカ', '移民法2002; BI-811', '非罹患証明', '2002'],
        ['', 'ナミビア', '移民管理法', '医学的禁止', '1993'],
        ['欧州', 'ロシア', '連邦法No.115-FZ', '国外退去事由', '2002'],
        ['', 'マルタ', '移民法Cap.217', '医学的入国不可', '\u2014'],
        ['南北アメリカ', '米国', 'INA §212(a)(1)(A)(i)', 'クラスA状態*', '1952\u2020'],
        ['', 'バルバドス', '移民法Cap.190', '医学的禁止', '\u2014'],
        ['', '米領バージン諸島', '米国連邦法', 'クラスA状態*', '\u2014'],
        ['南アジア', 'インド', '州レベル雇用法', '雇用制限', '各種'],
        ['', '香港特別行政区', '移民条例Cap.115', '医学的入国不可', '\u2014'],
    ]
    add_table_slide(prs,
        "表1：労働ビザ医療要件にハンセン病を明記している国（n=20）",
        t1h, t1d, [1.3, 1.5, 2.8, 2.0, 0.8], DARK_BLUE)

    # Slide 8: Table 2 - 15 countries without Hansen's (Japanese headers)
    t2h = ['国名', 'スクリーニング対象疾病', '特記事項', 'ハンセン病']
    t2d = [
        ['カナダ', '結核、梅毒', '移民医療検査（IME）', '含まれない'],
        ['オーストラリア', '結核、HIV、B/C型肝炎', 'パネル医師制度（Bupa）', '含まれない'],
        ['英国', '結核（6ヶ月超ビザ）', 'IOMパネル医師スクリーニング', '含まれない'],
        ['ニュージーランド', '結核、HIV、B型肝炎', 'Immigration NZパネル医師', '含まれない'],
        ['日本', '結核のみ（特定国籍）', 'JPETS入国前結核スクリーニング', '含まれない'],
        ['韓国', '結核、HIV、薬物', 'E-9ビザ医療検査', '含まれない'],
        ['シンガポール', '結核、HIV、梅毒、B型肝炎、マラリア', '外国人労働者医療検査', '含まれない'],
        ['イスラエル', '結核、HIV、B/C型肝炎', '入国医療要件', '含まれない'],
        ['ケニア', '結核、HIV', '労働許可医療証明書', '含まれない'],
        ['ナイジェリア', '結核、HIV、B型肝炎', '駐在員枠医療検査', '含まれない'],
        ['ガーナ', '結核、HIV、黄熱病', '移民医療証明書', '含まれない'],
        ['ブラジル', '黄熱病（予防接種）', '労働ビザ医療検査なし', '含まれない'],
        ['メキシコ', '指定なし', '標準医療検査なし', '含まれない'],
        ['アルゼンチン', '指定なし', '疾病特異的検査なし', '含まれない'],
        ['チリ', '結核、HIV（一部ビザ）', '一時居住者ビザ医療検査', '含まれない'],
    ]
    add_table_slide(prs,
        "表2：疾病特異的スクリーニングを行うがハンセン病を含まない国（n=15例）",
        t2h, t2d, [1.5, 2.5, 3.0, 1.4], DARK_GREEN)

    out = os.path.join(OUTPUT_DIR, "Scoping_Review_Figures_JA.pptx")
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == '__main__':
    create_pptx_ja()
